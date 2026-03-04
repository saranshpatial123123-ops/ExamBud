import fitz  # PyMuPDF
from pptx import Presentation

def extract_text_from_pdf(file_path: str) -> str:
    """Extracts text from a fully text-based PDF."""
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    return text

def extract_text_from_pptx(file_path: str) -> str:
    """Extracts text from a PowerPoint presentation."""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(file_path: str, file_type: str) -> str:
    """Routes the file to the correct extractor based on MIME type."""
    if file_type == "application/pdf":
        return extract_text_from_pdf(file_path)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
