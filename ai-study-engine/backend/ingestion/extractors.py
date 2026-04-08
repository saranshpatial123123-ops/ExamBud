import fitz  # PyMuPDF
from pptx import Presentation
import docx
import csv
import pandas as pd
from PIL import Image
import pytesseract

def extract_text_from_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    return text

def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_csv(file_path: str) -> str:
    df = pd.read_csv(file_path)
    return df.to_string()

def extract_text_from_image(file_path: str) -> str:
    img = Image.open(file_path)
    return pytesseract.image_to_string(img).strip()

def extract_text(file_path: str, file_type: str) -> str:
    """Routes the file to the correct extractor based on MIME type."""
    if file_type == "application/pdf":
        return extract_text_from_pdf(file_path)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_pptx(file_path)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_from_docx(file_path)
    elif file_type == "text/plain":
        return extract_text_from_txt(file_path)
    elif file_type == "text/csv":
        return extract_text_from_csv(file_path)
    elif file_type.startswith("image/"):
        return extract_text_from_image(file_path)
    else:
        raise ValueError(f"Unsupported document file type: {file_type}")
